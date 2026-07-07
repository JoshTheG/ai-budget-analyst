"""Dashboard snapshot, PowerPoint deck, and CLI output tests."""

import json
import subprocess
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
OPERATING = ROOT / "data" / "sample_prns_operating_budget.csv"
CAPITAL = ROOT / "data" / "sample_prns_capital_funds.csv"

sys.path.insert(0, str(ROOT))

from budget_analyst import dashboard  # noqa: E402


@pytest.fixture(scope="session", autouse=True)
def sample_data():
    if not OPERATING.exists() or not CAPITAL.exists():
        subprocess.run([sys.executable, str(ROOT / "data" / "make_sample.py")],
                       check=True)


def test_dashboard_snapshot_is_json_safe():
    snap = dashboard.compute_snapshot(str(CAPITAL))
    # must round-trip through strict JSON (NaN would raise)
    encoded = json.dumps(snap, default=str, allow_nan=False)
    decoded = json.loads(encoded)
    assert decoded["source"] == CAPITAL.name
    assert "fund_summary" in decoded["tables"]
    assert decoded["facts"]["total_budget"] > 0
    for table in decoded["tables"].values():
        assert set(table) == {"columns", "rows"}


def test_dashboard_state_caches_until_file_changes(tmp_path):
    src = tmp_path / "data.csv"
    src.write_text(OPERATING.read_text(encoding="utf-8"), encoding="utf-8")
    state = dashboard._State(str(src), None, None, "m")
    first = state.get()
    assert state.get() is first  # unchanged mtime -> cached object
    # append a row and bump mtime -> recompute
    with src.open("a", encoding="utf-8") as f:
        f.write('"FY 2025-26","New Division","Personal Services",'
                "1000000,900000,0,0\n")
    import os
    os.utime(src, (src.stat().st_atime, src.stat().st_mtime + 10))
    second = state.get()
    assert second is not first
    assert second["facts"]["total_budget"] == pytest.approx(
        first["facts"]["total_budget"] + 1_000_000)


def test_brief_cli_writes_deck(tmp_path):
    subprocess.run(
        [sys.executable, "-m", "budget_analyst", "brief", str(CAPITAL),
         "--out", str(tmp_path), "--mock"],
        cwd=ROOT, capture_output=True, text=True, check=True)
    deck_path = tmp_path / "budget_briefing.pptx"
    assert deck_path.exists()
    from pptx import Presentation
    prs = Presentation(str(deck_path))
    assert len(prs.slides) == 4
    text = "\n".join(sh.text_frame.text for slide in prs.slides
                     for sh in slide.shapes if sh.has_text_frame)
    assert "Budget Briefing" in text
    assert "Findings" in text


def test_analyze_capital_workbook_has_fund_sheet(tmp_path):
    subprocess.run(
        [sys.executable, "-m", "budget_analyst", "analyze", str(CAPITAL),
         "--out", str(tmp_path), "--mock"],
        cwd=ROOT, capture_output=True, text=True, check=True)
    from openpyxl import load_workbook
    wb = load_workbook(tmp_path / "budget_analysis.xlsx")
    assert "fund_summary" in wb.sheetnames
    assert "revenue_vs_target" in wb.sheetnames
