"""One-command PowerPoint briefing deck.

    python -m budget_analyst brief data/sample_prns_operating_budget.csv

Produces budget_briefing.pptx: title slide, headline figures, a native
budget-vs-actual chart, and findings/recommended actions. Like the memo,
every number is a pre-computed fact - the deck generator never does
arithmetic. Requires python-pptx (pip install python-pptx).
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

NAVY = (0x1F, 0x4E, 0x79)


def _fmt_money(v) -> str:
    try:
        return f"${v:,.0f}"
    except (TypeError, ValueError):
        return str(v)


def _bullets(facts: dict) -> tuple[list[str], list[str]]:
    """Finding and action bullets from computed facts only."""
    f, findings = facts, []
    if "overall_pct_spent" in f:
        findings.append(
            f"{f['overall_pct_spent']}% of the {_fmt_money(f['total_budget'])} budget "
            f"spent ({_fmt_money(f['total_actual'])}); net variance "
            f"{_fmt_money(f['total_variance'])}.")
    if "overall_pct_committed" in f:
        findings.append(
            f"Including {_fmt_money(f['total_encumbrance'])} encumbered, "
            f"{f['overall_pct_committed']}% of budget is committed; "
            f"{_fmt_money(f['total_available'])} remains available.")
    if "largest_overrun_entity" in f:
        findings.append(
            f"Highest spend vs. budget: {f['largest_overrun_entity']} "
            f"({f['largest_overrun_pct']:+.1f}%); largest underspend: "
            f"{f['largest_underspend_entity']} ({f['largest_underspend_pct']:+.1f}%).")
    if "revenue_attainment_pct" in f:
        findings.append(
            f"Revenue at {f['revenue_attainment_pct']}% of the "
            f"{_fmt_money(f['total_revenue_target'])} target; weakest: "
            f"{f['weakest_revenue_entity']} "
            f"({f['weakest_revenue_attainment_pct']:.1f}%).")
    if "tightest_fund" in f:
        findings.append(
            f"Across {f['n_funds']} funds, tightest available balance is "
            f"{f['tightest_fund']} ({_fmt_money(f['tightest_fund_available'])}).")
    if "forecast_next_period" in f:
        findings.append(
            f"Next-period projection: {_fmt_money(f['forecast_next_period'])} "
            f"({f['forecast_method']}, R² = {f['forecast_r_squared']}).")

    actions = ["Review flagged outlier units with program managers before close."
               if f.get("anomaly_count")
               else "No statistical outliers flagged; continue standard monitoring.",
               "Reconcile top variance lines against encumbrances and timing effects.",
               "Re-run when the next actuals post; the deck rebuilds in one command."]
    return findings, actions


def write_deck(result: dict, source_name: str, out_dir: str) -> str:
    """Build budget_briefing.pptx from the analysis result."""
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for the brief command: "
            "pip install python-pptx") from exc

    facts, tables = result["facts"], result["tables"]
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_title(slide, text: str, size: int = 28, top: float = 0.35):
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(0.9))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size, p.font.bold = Pt(size), True
        p.font.color.rgb = RGBColor(*NAVY)
        return box

    # --- 1: title
    s = prs.slides.add_slide(blank)
    add_title(s, "Budget Briefing", 44, 2.4)
    sub = s.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(12), Inches(1.4))
    tf = sub.text_frame
    tf.text = f"Source: {source_name}"
    p = tf.add_paragraph()
    p.text = (f"Period: {facts.get('latest_period', 'n/a')}  |  "
              f"Prepared {date.today():%B %d, %Y}  |  AI Budget Analyst")
    for para in tf.paragraphs:
        para.font.size = Pt(16)

    # --- 2: headline figures
    s = prs.slides.add_slide(blank)
    add_title(s, "Headline Figures")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    pairs = [("Total budget", _fmt_money(facts.get("total_budget"))),
             ("Actual expenditures", _fmt_money(facts.get("total_actual"))),
             ("% of budget spent", f"{facts.get('overall_pct_spent', 'n/a')}%")]
    if "total_available" in facts:
        pairs += [("Encumbered", _fmt_money(facts["total_encumbrance"])),
                  ("Available balance", _fmt_money(facts["total_available"]))]
    if "revenue_attainment_pct" in facts:
        pairs.append(("Revenue attainment",
                      f"{facts['revenue_attainment_pct']}% of "
                      f"{_fmt_money(facts['total_revenue_target'])}"))
    for i, (label, value) in enumerate(pairs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{label}:  {value}"
        p.font.size = Pt(24)
        p.space_after = Pt(14)

    # --- 3: chart
    if "variance_by_entity" in tables:
        var = tables["variance_by_entity"].head(10)
        ent = facts.get("entity_column", var.columns[0])
        s = prs.slides.add_slide(blank)
        add_title(s, f"Budget vs. Actual by {ent}")
        data = CategoryChartData()
        data.categories = [str(v)[:28] for v in var[ent]]
        data.add_series("Budget", [float(v) for v in var["budget"]])
        data.add_series("Actual", [float(v) for v in var["actual"]])
        chart = s.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.3),
            Inches(12.1), Inches(5.7), data).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    # --- 4: findings & actions
    s = prs.slides.add_slide(blank)
    add_title(s, "Findings & Recommended Actions")
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    findings, actions = _bullets(facts)
    first = True
    for header, items in (("Findings", findings), ("Recommended actions", actions)):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = header
        p.font.size, p.font.bold = Pt(20), True
        p.font.color.rgb = RGBColor(*NAVY)
        for item in items:
            q = tf.add_paragraph()
            q.text = f"• {item}"
            q.font.size = Pt(15)
            q.space_after = Pt(6)

    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    path = out / "budget_briefing.pptx"
    prs.save(path)
    return str(path)
